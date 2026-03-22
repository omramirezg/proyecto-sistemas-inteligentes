from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parent.parent
DOCS = ROOT / "docs"
OUT = DOCS / "Presentacion_Sustentacion_LLM_Multimodal_UNAL_v5.pptx"
UNAL_LOGO = DOCS / "unal_logo_template.png"


COLOR_FONDO = RGBColor(247, 243, 236)
COLOR_NAVY = RGBColor(24, 44, 62)
COLOR_AZUL = RGBColor(34, 78, 122)
COLOR_TEAL = RGBColor(39, 110, 104)
COLOR_VERDE = RGBColor(90, 122, 60)
COLOR_AMBAR = RGBColor(176, 125, 62)
COLOR_ROJO = RGBColor(150, 67, 54)
COLOR_GRIS = RGBColor(86, 90, 94)
COLOR_CLARO = RGBColor(253, 251, 247)
COLOR_BLANCO = RGBColor(255, 255, 255)
COLOR_CODIGO = RGBColor(30, 34, 41)

FONT_TITULO = "Aptos Display"
FONT_TEXTO = "Aptos"
FONT_CODIGO = "Consolas"


def limpiar_slide(slide):
    for shape in list(slide.shapes):
        sp = shape._element
        sp.getparent().remove(sp)


def fondo(slide, color=COLOR_FONDO):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def barra_superior(slide, titulo, subtitulo=None):
    bar = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, Inches(13.333), Inches(0.75)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = COLOR_NAVY
    bar.line.fill.background()

    tx = slide.shapes.add_textbox(Inches(0.55), Inches(0.18), Inches(8.8), Inches(0.34))
    tf = tx.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = titulo
    r.font.name = FONT_TITULO
    r.font.size = Pt(24)
    r.font.bold = True
    r.font.color.rgb = COLOR_BLANCO

    if subtitulo:
        st = slide.shapes.add_textbox(Inches(9.45), Inches(0.2), Inches(3.1), Inches(0.28))
        tf = st.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.RIGHT
        r = p.add_run()
        r.text = subtitulo
        r.font.name = FONT_TEXTO
        r.font.size = Pt(11)
        r.font.color.rgb = RGBColor(220, 226, 233)


def pie_unal(slide, texto="Facultad de Minas | Sede Medellin"):
    pie = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, Inches(6.9), Inches(13.333), Inches(0.6)
    )
    pie.fill.solid()
    pie.fill.fore_color.rgb = COLOR_NAVY
    pie.line.fill.background()

    txt = slide.shapes.add_textbox(Inches(0.7), Inches(7.0), Inches(5.5), Inches(0.24))
    tf = txt.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = texto
    r.font.name = FONT_TEXTO
    r.font.size = Pt(12)
    r.font.bold = True
    r.font.color.rgb = COLOR_BLANCO


def caja_texto(slide, x, y, w, h, titulo, cuerpo, color_titulo=COLOR_NAVY):
    box = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, y, w, h)
    box.fill.solid()
    box.fill.fore_color.rgb = COLOR_CLARO
    box.line.color.rgb = RGBColor(224, 220, 210)

    title = slide.shapes.add_textbox(x + Inches(0.18), y + Inches(0.14), w - Inches(0.35), Inches(0.28))
    tf = title.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = titulo
    r.font.name = FONT_TITULO
    r.font.bold = True
    r.font.size = Pt(17)
    r.font.color.rgb = color_titulo

    body = slide.shapes.add_textbox(x + Inches(0.2), y + Inches(0.52), w - Inches(0.35), h - Inches(0.65))
    tf = body.text_frame
    tf.word_wrap = True
    tf.margin_bottom = 0
    tf.margin_top = 0
    tf.margin_left = 0
    tf.margin_right = 0
    for i, line in enumerate(cuerpo):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.level = 0
        p.font.name = FONT_TEXTO
        p.font.size = Pt(15)
        p.font.color.rgb = COLOR_GRIS
        p.bullet = True


def tarjeta_kpi(slide, x, y, w, h, titulo, valor, color):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()

    tf = shape.text_frame
    tf.clear()
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p1 = tf.paragraphs[0]
    p1.alignment = PP_ALIGN.CENTER
    r1 = p1.add_run()
    r1.text = titulo
    r1.font.name = FONT_TEXTO
    r1.font.size = Pt(12)
    r1.font.bold = True
    r1.font.color.rgb = COLOR_BLANCO

    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run()
    r2.text = valor
    r2.font.name = FONT_TITULO
    r2.font.size = Pt(22)
    r2.font.bold = True
    r2.font.color.rgb = COLOR_BLANCO


def codigo_box(slide, x, y, w, h, titulo, codigo):
    titulo_box = slide.shapes.add_textbox(x, y - Inches(0.32), w, Inches(0.25))
    tf = titulo_box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = titulo
    r.font.name = FONT_TITULO
    r.font.size = Pt(16)
    r.font.bold = True
    r.font.color.rgb = COLOR_NAVY

    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLOR_CODIGO
    shape.line.fill.background()

    text = slide.shapes.add_textbox(x + Inches(0.18), y + Inches(0.18), w - Inches(0.3), h - Inches(0.28))
    tf = text.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.clear()
    for i, line in enumerate(codigo.splitlines()):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.name = FONT_CODIGO
        p.font.size = Pt(11)
        p.font.color.rgb = RGBColor(236, 240, 241)
        p.space_after = 0
        p.space_before = 0


def nota_presentador(slide, titulo, texto):
    box = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(9.3), Inches(6.45), Inches(3.45), Inches(0.7)
    )
    box.fill.solid()
    box.fill.fore_color.rgb = COLOR_BLANCO
    box.line.color.rgb = RGBColor(221, 214, 199)
    tf = box.text_frame
    tf.clear()
    p1 = tf.paragraphs[0]
    r1 = p1.add_run()
    r1.text = titulo
    r1.font.name = FONT_TITULO
    r1.font.bold = True
    r1.font.size = Pt(11)
    r1.font.color.rgb = COLOR_AZUL
    p2 = tf.add_paragraph()
    r2 = p2.add_run()
    r2.text = texto
    r2.font.name = FONT_TEXTO
    r2.font.size = Pt(10)
    r2.font.color.rgb = COLOR_GRIS


prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)


def slide_titulo():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    limpiar_slide(slide)
    fondo(slide, COLOR_FONDO)

    acento = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.0), Inches(0.0), Inches(0.42), Inches(7.5)
    )
    acento.fill.solid()
    acento.fill.fore_color.rgb = COLOR_TEAL
    acento.line.fill.background()

    banda = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.78), Inches(0.72), Inches(4.75), Inches(0.44)
    )
    banda.fill.solid()
    banda.fill.fore_color.rgb = RGBColor(224, 234, 228)
    banda.line.fill.background()
    tf = banda.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = "Universidad Nacional de Colombia"
    r.font.name = FONT_TEXTO
    r.font.size = Pt(15)
    r.font.bold = True
    r.font.color.rgb = COLOR_NAVY

    title = slide.shapes.add_textbox(Inches(0.82), Inches(1.45), Inches(6.1), Inches(2.6))
    tf = title.text_frame
    tf.word_wrap = True
    tf.clear()
    p1 = tf.paragraphs[0]
    r1 = p1.add_run()
    r1.text = "Sistema de alertamiento prescriptivo multimodal"
    r1.font.name = FONT_TITULO
    r1.font.size = Pt(24)
    r1.font.bold = True
    r1.font.color.rgb = COLOR_NAVY

    p1b = tf.add_paragraph()
    r1b = p1b.add_run()
    r1b.text = "para peletizacion industrial"
    r1b.font.name = FONT_TITULO
    r1b.font.size = Pt(24)
    r1b.font.bold = True
    r1b.font.color.rgb = COLOR_NAVY

    p2 = tf.add_paragraph()
    r2 = p2.add_run()
    r2.text = "Integracion de reglas deterministas, Gemini en Vertex AI, Nano Banana, Telegram y TTS"
    r2.font.name = FONT_TEXTO
    r2.font.size = Pt(15)
    r2.font.color.rgb = COLOR_GRIS

    tarjeta_kpi(slide, Inches(0.85), Inches(4.58), Inches(2.05), Inches(1.05), "Canal operativo", "Telegram", COLOR_AZUL)
    tarjeta_kpi(slide, Inches(3.1), Inches(4.58), Inches(2.05), Inches(1.05), "IA visual", "Nano Banana", COLOR_TEAL)
    tarjeta_kpi(slide, Inches(5.35), Inches(4.58), Inches(2.05), Inches(1.05), "Prescripcion", "Gemini + TTS", COLOR_VERDE)

    box = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(7.7), Inches(0.78), Inches(4.85), Inches(5.95))
    box.fill.solid()
    box.fill.fore_color.rgb = COLOR_BLANCO
    box.line.color.rgb = RGBColor(225, 216, 201)

    img = DOCS / "DiagramaDelProceso.png"
    if img.exists():
        slide.shapes.add_picture(str(img), Inches(8.0), Inches(1.0), width=Inches(4.25), height=Inches(5.5))

    footer = slide.shapes.add_textbox(Inches(0.84), Inches(6.02), Inches(6.2), Inches(0.34))
    tf = footer.text_frame
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = "Sustentacion enfocada en codigo, funcionalidad operativa y uso real de LLMs multimodales"
    r.font.name = FONT_TEXTO
    r.font.size = Pt(12)
    r.font.color.rgb = COLOR_GRIS

    pie_unal(slide, "Facultad de Minas | Sede Medellin")


def slide_problema():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    limpiar_slide(slide)
    fondo(slide)
    barra_superior(slide, "1. Problema Industrial", "Contexto y necesidad")

    caja_texto(
        slide, Inches(0.7), Inches(1.2), Inches(4.0), Inches(2.0),
        "Desafio en planta",
        [
            "La peletizacion depende de temperatura, presion y carga en tiempo real.",
            "El operario trabaja en un entorno ruidoso, con alta demanda de reaccion y poca tolerancia al error.",
            "Ver solo numeros no siempre permite entender tendencia, severidad ni accion inmediata.",
        ],
    )
    caja_texto(
        slide, Inches(4.95), Inches(1.2), Inches(3.7), Inches(2.0),
        "Riesgos operativos",
        [
            "Fatiga de alarmas.",
            "Retrasos de respuesta.",
            "Producto fuera de especificacion.",
            "Sobrecarga del equipo y uso ineficiente del vapor.",
        ],
        color_titulo=COLOR_ROJO,
    )
    caja_texto(
        slide, Inches(8.95), Inches(1.2), Inches(3.65), Inches(2.0),
        "Objetivo del proyecto",
        [
            "Convertir telemetria en apoyo operacional claro, prescriptivo y multimodal.",
            "Entregar informacion util por Telegram para operario y gerente.",
        ],
        color_titulo=COLOR_TEAL,
    )

    tarjeta_kpi(slide, Inches(1.0), Inches(4.0), Inches(2.3), Inches(1.1), "Lectura simulada", "1 min", COLOR_AZUL)
    tarjeta_kpi(slide, Inches(3.55), Inches(4.0), Inches(2.3), Inches(1.1), "Worker backend", "24/7", COLOR_TEAL)
    tarjeta_kpi(slide, Inches(6.1), Inches(4.0), Inches(2.3), Inches(1.1), "Canal de accion", "Telegram", COLOR_VERDE)
    tarjeta_kpi(slide, Inches(8.65), Inches(4.0), Inches(3.0), Inches(1.1), "Salida multimodal", "texto + imagen + audio", COLOR_AMBAR)

    nota_presentador(slide, "Mensaje clave", "No es un chatbot; es soporte operacional de planta en tiempo real.")
    pie_unal(slide, "Facultad de Minas | Sede Medellin")


def slide_solucion():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    limpiar_slide(slide)
    fondo(slide)
    barra_superior(slide, "2. Propuesta de Solucion", "Que construimos")

    caja_texto(
        slide, Inches(0.8), Inches(1.15), Inches(5.75), Inches(4.6),
        "Sistema propuesto",
        [
            "Backend perpetuo que consume telemetria simulada desde CSV como si fuera un sensor IoT.",
            "Motor determinista con EMA, bandas operativas, persistencia de alarma y antispam.",
            "Panel tecnico de proceso, prescripcion de Maria, audio TTS, fichas IA, PDF y dashboard.",
            "Interaccion por Telegram con botones para operario, gerente y feedback humano.",
        ],
    )

    flow = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(6.9), Inches(1.15), Inches(5.55), Inches(4.6))
    flow.fill.solid()
    flow.fill.fore_color.rgb = COLOR_BLANCO
    flow.line.color.rgb = RGBColor(223, 218, 208)
    tf = flow.text_frame
    tf.clear()
    pasos = [
        "1. CSV de telemetria planta 001",
        "2. DataLoader y maestros",
        "3. Motor de reglas + EMA",
        "4. Panel visual del proceso",
        "5. Gemini / Maria + Nano Banana",
        "6. Google TTS + Telegram",
        "7. Historial, feedback, dashboard y PDF",
    ]
    for i, paso in enumerate(pasos):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = paso
        p.font.name = FONT_TEXTO
        p.font.size = Pt(18 if i == 0 else 16)
        p.font.color.rgb = COLOR_NAVY if i == 0 else COLOR_GRIS
        p.bullet = False
        p.space_after = Pt(12 if i == 0 else 10)
    nota_presentador(slide, "Mensaje clave", "La seguridad de disparo es determinista; la capa generativa agrega interpretacion y presentacion.")
    pie_unal(slide, "Facultad de Minas | Sede Medellin")


def slide_arquitectura():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    limpiar_slide(slide)
    fondo(slide)
    barra_superior(slide, "3. Arquitectura del Sistema", "Pipeline de extremo a extremo")

    img = DOCS / "DiagramaDelProceso.png"
    marco = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.7), Inches(1.1), Inches(8.1), Inches(5.7))
    marco.fill.solid()
    marco.fill.fore_color.rgb = COLOR_BLANCO
    marco.line.color.rgb = RGBColor(220, 214, 202)
    if img.exists():
        slide.shapes.add_picture(str(img), Inches(0.95), Inches(1.3), width=Inches(7.6), height=Inches(5.25))

    caja_texto(
        slide, Inches(9.1), Inches(1.15), Inches(3.55), Inches(5.4),
        "Capas principales",
        [
            "Ingesta: telemetria y maestros.",
            "Control: EMA, bandas, persistencia, severidad.",
            "Multimodalidad: panel tecnico, Gemini, Nano Banana.",
            "Entrega: Telegram, audio, fichas, PDF, dashboard.",
            "Trazabilidad: historial de alertas y feedback del operario.",
        ],
        color_titulo=COLOR_AZUL,
    )
    nota_presentador(slide, "Mensaje clave", "El repositorio esta desacoplado: datos, reglas, IA, TTS, Telegram y reportes.")
    pie_unal(slide, "Facultad de Minas | Sede Medellin")


def slide_funcionalidad():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    limpiar_slide(slide)
    fondo(slide)
    barra_superior(slide, "4. Funcionalidad Demostrada", "Que ve el usuario en operacion")

    caja_texto(
        slide, Inches(0.8), Inches(1.15), Inches(3.8), Inches(4.9),
        "Operario en Telegram",
        [
            "Mensaje tecnico con estado global, severidad, salud y causa probable.",
            "Panel visual del proceso.",
            "Audio prescriptivo para entorno ruidoso.",
            "Botones: PDF, Dashboard, Ficha Operario, Ficha Gerencial, Explicar evento.",
            "Feedback del operario: util, falso positivo, mantenimiento.",
        ],
        color_titulo=COLOR_TEAL,
    )
    caja_texto(
        slide, Inches(4.85), Inches(1.15), Inches(3.8), Inches(4.9),
        "Gerencia y supervision",
        [
            "PDF ejecutivo bajo demanda.",
            "Dashboard HTML local.",
            "Ficha gerencial visual generada por IA.",
            "Resumen de alertas, feedback y estabilidad del proceso.",
        ],
        color_titulo=COLOR_AMBAR,
    )
    caja_texto(
        slide, Inches(8.9), Inches(1.15), Inches(3.8), Inches(4.9),
        "Valor practico",
        [
            "No solo detecta desviaciones.",
            "Explica tendencia y pronostico.",
            "Convierte datos crudos en decisiones legibles.",
            "Reduce saturacion al separar salida operativa de salida ejecutiva.",
        ],
        color_titulo=COLOR_ROJO,
    )
    nota_presentador(slide, "Mensaje clave", "La demo debe mostrar una lectura normal, una baja y una alta con reaccion distinta.")
    pie_unal(slide, "Facultad de Minas | Sede Medellin")


def slide_evidencia_visual():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    limpiar_slide(slide)
    fondo(slide)
    barra_superior(slide, "5. Evidencia Visual en Telegram", "Panel tecnico y fichas IA")

    assets = DOCS / "presentacion_assets"
    panel = assets / "panel_real_demo.jpg"
    ficha = assets / "ficha_real_demo.jpg"
    telegram = assets / "telegram_real_demo.jpg"

    if not panel.exists():
        panel = assets / "panel_telegram_demo.png"
    if not ficha.exists():
        ficha = assets / "ficha_ia_mock_demo.png"
    if not telegram.exists():
        telegram = assets / "telegram_mock_demo.png"

    if telegram.exists():
        marco_tg = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.75), Inches(1.15), Inches(3.35), Inches(5.45))
        marco_tg.fill.solid()
        marco_tg.fill.fore_color.rgb = COLOR_BLANCO
        marco_tg.line.color.rgb = RGBColor(222, 214, 201)
        slide.shapes.add_picture(str(telegram), Inches(0.95), Inches(1.35), width=Inches(2.95), height=Inches(4.95))

    if panel.exists():
        marco_panel = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(4.35), Inches(1.15), Inches(4.05), Inches(2.55))
        marco_panel.fill.solid()
        marco_panel.fill.fore_color.rgb = COLOR_BLANCO
        marco_panel.line.color.rgb = RGBColor(222, 214, 201)
        slide.shapes.add_picture(str(panel), Inches(4.55), Inches(1.35), width=Inches(3.65), height=Inches(2.15))

    if ficha.exists():
        marco_ficha = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(8.7), Inches(1.15), Inches(3.9), Inches(5.1))
        marco_ficha.fill.solid()
        marco_ficha.fill.fore_color.rgb = COLOR_BLANCO
        marco_ficha.line.color.rgb = RGBColor(222, 214, 201)
        slide.shapes.add_picture(str(ficha), Inches(8.92), Inches(1.35), width=Inches(3.45), height=Inches(4.6))

    caja_texto(
        slide, Inches(4.35), Inches(4.0), Inches(4.05), Inches(2.0),
        "Que aporta esta evidencia",
        [
            "Se muestra el chat real de Telegram, el panel tecnico real y una ficha IA real generada para el proyecto.",
            "Esto permite sustentar funcionalidad operativa y no solo diseno conceptual.",
            "La diapositiva deja visible la cadena completa: mensaje, visual analitica y salida generativa.",
        ],
        color_titulo=COLOR_TEAL,
    )
    nota_presentador(slide, "Mensaje clave", "Aqui conviene parar y decir: esto es lo que realmente recibe el usuario final.")
    pie_unal(slide, "Facultad de Minas | Sede Medellin")


def slide_llms():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    limpiar_slide(slide)
    fondo(slide)
    barra_superior(slide, "6. Uso de LLMs y Multimodalidad", "Donde interviene la IA")

    caja_texto(
        slide, Inches(0.75), Inches(1.15), Inches(3.9), Inches(2.2),
        "Gemini / Maria",
        [
            "Recibe contexto estructurado + panel tecnico.",
            "Genera prescripcion breve para audio y chat.",
            "Tambien explica el evento bajo demanda.",
        ],
        color_titulo=COLOR_AZUL,
    )
    caja_texto(
        slide, Inches(4.75), Inches(1.15), Inches(3.9), Inches(2.2),
        "Nano Banana",
        [
            "Usa la imagen tecnica como referencia multimodal.",
            "Genera fichas visuales nuevas para operario y gerencia.",
            "No replica la grafica: la transforma en una pieza ejecutiva.",
        ],
        color_titulo=COLOR_TEAL,
    )
    caja_texto(
        slide, Inches(8.75), Inches(1.15), Inches(3.85), Inches(2.2),
        "Google TTS",
        [
            "Convierte la prescripcion en audio reproducible en Telegram.",
            "Permite asistencia operativa en ambiente ruidoso.",
        ],
        color_titulo=COLOR_VERDE,
    )

    codigo = (
        "prescripcion_maria = self._generar_prescripcion_maria(..., imagen_bytes=imagen_bytes)\n"
        "audio_texto = prescripcion_maria\n"
        "audio_bytes = self.tts.sintetizar(audio_texto)\n\n"
        "imagen_ia, _ = self.nano_banana.generar_ficha_visual(\n"
        "    prompt_texto=prompt,\n"
        "    imagen_referencia_bytes=self._ultimo_panel_bytes,\n"
        ")"
    )
    codigo_box(slide, Inches(0.9), Inches(4.1), Inches(7.0), Inches(1.85), "Fragmento representativo", codigo)
    caja_texto(
        slide, Inches(8.2), Inches(4.05), Inches(4.3), Inches(2.0),
        "Punto academico",
        [
            "La multimodalidad no es decorativa.",
            "Se combina texto, imagen tecnica y generacion visual.",
            "La capa generativa opera sobre una base determinista y trazable.",
        ],
        color_titulo=COLOR_ROJO,
    )
    nota_presentador(slide, "Mensaje clave", "Aqui es donde justificamos que si usamos LLMs multimodales de forma funcional.")
    pie_unal(slide, "Facultad de Minas | Sede Medellin")


def slide_codigo_motor():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    limpiar_slide(slide)
    fondo(slide)
    barra_superior(slide, "7. Codigo Clave: Motor y Persistencia", "Mostrar en Antigravity durante la exposicion")

    codigo = (
        "if nuevo_estado == estado_actual:\n"
        "    persistencia = 0\n"
        "    return None\n\n"
        "if nuevo_estado != 'NORMAL':\n"
        "    persistencia += 1\n"
        "    if persistencia < self.confirmacion_alerta:\n"
        "        return None\n\n"
        "if nuevo_estado == 'NORMAL':\n"
        "    persistencia += 1\n"
        "    if persistencia < self.confirmacion_retorno:\n"
        "        return None"
    )
    codigo_box(slide, Inches(0.8), Inches(1.55), Inches(6.2), Inches(3.8), "motor_reglas.py", codigo)
    caja_texto(
        slide, Inches(7.35), Inches(1.55), Inches(5.2), Inches(3.8),
        "Que explico aqui",
        [
            "No se dispara una alarma por una sola lectura atipica.",
            "La desviacion debe persistir para confirmarse.",
            "Esto reduce ruido operacional y fatiga de alarmas.",
            "Es una decision ingenieril importante en la tesis.",
        ],
        color_titulo=COLOR_AZUL,
    )
    nota_presentador(slide, "En vivo", "Abrir [motor_reglas.py](C:/Users/USUARIO/Desktop/Maestria/sistemasinteligentes/proyecto/src/motor_reglas.py) y mostrar la confirmacion por persistencia.")
    pie_unal(slide, "Facultad de Minas | Sede Medellin")


def slide_codigo_main():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    limpiar_slide(slide)
    fondo(slide)
    barra_superior(slide, "8. Codigo Clave: Orquestacion Multimodal", "Worker principal")

    codigo = (
        "imagen_bytes = self.graficas.generar_panel_multimodal_telegram(...)\n\n"
        "prescripcion_maria = self._generar_prescripcion_maria(\n"
        "    lectura=lectura,\n"
        "    ...,\n"
        "    imagen_bytes=imagen_bytes,\n"
        ")\n\n"
        "texto += f\"\\n<b>Maria:</b> {html.escape(prescripcion_maria)}\"\n"
        "audio_bytes = self.tts.sintetizar(prescripcion_maria)"
    )
    codigo_box(slide, Inches(0.8), Inches(1.55), Inches(6.3), Inches(3.8), "main.py", codigo)
    caja_texto(
        slide, Inches(7.4), Inches(1.55), Inches(5.1), Inches(3.8),
        "Que explico aqui",
        [
            "El worker no solo calcula estados.",
            "Construye un panel tecnico, llama a Maria y convierte su salida en audio.",
            "La respuesta generativa se apoya en la imagen del proceso y en contexto numerico.",
            "Si la IA falla, el sistema conserva un fallback determinista.",
        ],
        color_titulo=COLOR_TEAL,
    )
    nota_presentador(slide, "En vivo", "Abrir [main.py](C:/Users/USUARIO/Desktop/Maestria/sistemasinteligentes/proyecto/src/main.py) y mostrar donde entra Gemini y donde entra Nano Banana.")
    pie_unal(slide, "Facultad de Minas | Sede Medellin")


def slide_codigo_telegram():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    limpiar_slide(slide)
    fondo(slide)
    barra_superior(slide, "9. Codigo Clave: Interaccion por Telegram", "Botones, callbacks y feedback")

    codigo = (
        "filas.append([\n"
        "    InlineKeyboardButton('Ficha Operario', callback_data='solicitar_ficha_operario'),\n"
        "    InlineKeyboardButton('Ficha Gerencial', callback_data='solicitar_ficha_gerencial'),\n"
        "])\n"
        "filas.append([\n"
        "    InlineKeyboardButton('Explicar evento', callback_data='explicar_evento'),\n"
        "])\n\n"
        "elif callback_query.data == 'explicar_evento':\n"
        "    eventos['explicar_evento'].append(chat_id)"
    )
    codigo_box(slide, Inches(0.8), Inches(1.5), Inches(6.35), Inches(3.9), "telegram_bot.py", codigo)
    caja_texto(
        slide, Inches(7.45), Inches(1.5), Inches(5.0), Inches(3.9),
        "Que explico aqui",
        [
            "Telegram no es solo un canal de salida.",
            "Tambien es una interfaz operativa con botones y feedback.",
            "Permite pedir PDF, dashboard, fichas IA y explicaciones del evento.",
            "El operario puede validar si una alerta fue util o si requiere mantenimiento.",
        ],
        color_titulo=COLOR_AMBAR,
    )
    nota_presentador(slide, "En vivo", "Abrir [telegram_bot.py](C:/Users/USUARIO/Desktop/Maestria/sistemasinteligentes/proyecto/src/telegram_bot.py) y luego mostrar el chat del bot.")
    pie_unal(slide, "Facultad de Minas | Sede Medellin")


def slide_demo():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    limpiar_slide(slide)
    fondo(slide)
    barra_superior(slide, "10. Demostracion Sugerida", "Guion de 10 minutos")

    caja_texto(
        slide, Inches(0.8), Inches(1.15), Inches(3.8), Inches(5.0),
        "Escenario 1",
        [
            "Lecturas 1 a 5 en banda normal.",
            "Mostrar estado estable, ficha de continuidad y mensaje de Maria.",
            "Explicar que no hay alarmas por ruido.",
        ],
        color_titulo=COLOR_VERDE,
    )
    caja_texto(
        slide, Inches(4.8), Inches(1.15), Inches(3.8), Inches(5.0),
        "Escenario 2",
        [
            "Lecturas 6 a 8 con temperatura y presion bajas.",
            "Mostrar persistencia, severidad y recomendacion operativa.",
            "Pedir Ficha Operario y Explicar evento.",
        ],
        color_titulo=COLOR_AMBAR,
    )
    caja_texto(
        slide, Inches(8.8), Inches(1.15), Inches(3.8), Inches(5.0),
        "Escenario 3",
        [
            "Lecturas 13 a 14 con temperatura y presion altas.",
            "Mostrar Ficha Gerencial y PDF bajo demanda.",
            "Cerrar con el dashboard y el feedback del operario.",
        ],
        color_titulo=COLOR_ROJO,
    )
    nota_presentador(slide, "Consejo", "Alternar slides y Antigravity: no meter demasiado codigo en PowerPoint.")
    pie_unal(slide, "Facultad de Minas | Sede Medellin")


def slide_aporte():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    limpiar_slide(slide)
    fondo(slide)
    barra_superior(slide, "11. Aporte Academico y Tecnico", "Por que este proyecto destaca")

    caja_texto(
        slide, Inches(0.85), Inches(1.2), Inches(5.8), Inches(4.9),
        "Aportes",
        [
            "Integra capa determinista y capa generativa sin delegar seguridad al LLM.",
            "Usa multimodalidad de forma funcional: imagen tecnica, interpretacion LLM, audio y fichas visuales.",
            "Incluye operacion en tiempo real, trazabilidad, feedback humano y salidas diferenciadas por audiencia.",
            "Convierte un problema industrial en una experiencia operativa explicable.",
        ],
        color_titulo=COLOR_AZUL,
    )
    caja_texto(
        slide, Inches(6.95), Inches(1.2), Inches(5.4), Inches(4.9),
        "Mensaje final al profesor",
        [
            "No se trata solo de usar Gemini.",
            "Se trata de diseñar un sistema completo donde la IA agrega valor real a la toma de decisiones.",
            "El proyecto muestra codigo, arquitectura, funcionalidad y multimodalidad aplicada a un contexto industrial.",
        ],
        color_titulo=COLOR_TEAL,
    )
    nota_presentador(slide, "Cierre", "El diferenciador no es una sola API; es la orquestacion de datos, reglas, LLMs y canal operativo.")
    pie_unal(slide, "Facultad de Minas | Sede Medellin")


def slide_cierre():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    limpiar_slide(slide)
    fondo(slide, COLOR_NAVY)

    tx = slide.shapes.add_textbox(Inches(1.05), Inches(1.3), Inches(11.0), Inches(1.3))
    tf = tx.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = "Gracias"
    r.font.name = FONT_TITULO
    r.font.size = Pt(32)
    r.font.bold = True
    r.font.color.rgb = COLOR_BLANCO

    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run()
    r2.text = "Sistema de alertamiento prescriptivo multimodal para peletizacion industrial"
    r2.font.name = FONT_TEXTO
    r2.font.size = Pt(18)
    r2.font.color.rgb = RGBColor(210, 220, 231)

    tarjeta_kpi(slide, Inches(2.0), Inches(3.2), Inches(2.6), Inches(1.2), "Mostrar en vivo", "Telegram", COLOR_AZUL)
    tarjeta_kpi(slide, Inches(5.35), Inches(3.2), Inches(2.6), Inches(1.2), "Mostrar codigo", "Antigravity", COLOR_TEAL)
    tarjeta_kpi(slide, Inches(8.7), Inches(3.2), Inches(2.6), Inches(1.2), "Cerrar con", "Ficha IA + PDF", COLOR_AMBAR)

    txt = slide.shapes.add_textbox(Inches(2.3), Inches(5.2), Inches(8.6), Inches(0.5))
    tf = txt.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = "Preguntas y demostracion operativa"
    r.font.name = FONT_TITULO
    r.font.size = Pt(20)
    r.font.bold = True
    r.font.color.rgb = COLOR_BLANCO

    pie_unal(slide, "Facultad de Minas | Sede Medellin")


slide_titulo()
slide_problema()
slide_solucion()
slide_arquitectura()
slide_funcionalidad()
slide_evidencia_visual()
slide_llms()
slide_codigo_motor()
slide_codigo_main()
slide_codigo_telegram()
slide_demo()
slide_aporte()
slide_cierre()

prs.save(str(OUT))
print(OUT)
